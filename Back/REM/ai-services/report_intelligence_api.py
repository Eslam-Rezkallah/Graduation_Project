import logging
import os
import re
from contextlib import asynccontextmanager
from pathlib import Path
from typing import Optional
from urllib.parse import urlparse, unquote
from zipfile import ZipFile
from xml.etree import ElementTree as ET

import pandas as pd
from fastapi import FastAPI, HTTPException, Query
from fastapi.middleware.cors import CORSMiddleware
from motor.motor_asyncio import AsyncIOMotorClient
from pydantic import BaseModel

try:
    from PyPDF2 import PdfReader
except ImportError:
    PdfReader = None

try:
    import fitz
except ImportError:
    fitz = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from openpyxl import load_workbook
except ImportError:
    load_workbook = None


logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(levelname)s] %(message)s")
log = logging.getLogger(__name__)

MONGO_URI = os.getenv("MONGO_URI", "mongodb://localhost:27017")
MONGO_DB = os.getenv("MONGO_DB", "REM")
MONGO_COLLECTION = os.getenv("MONGO_COLLECTION", "files")
FILE_PATH_FIELD = os.getenv("FILE_PATH_FIELD", "url")
FILE_BASE_DIR = Path(os.getenv("FILE_BASE_DIR", ".")).resolve()

mongo_client: AsyncIOMotorClient | None = None
_store: dict = {"tasks": [], "skills": [], "scanned": False, "file_count": 0}

NS = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}


class TaskRecord(BaseModel):
    employee: str
    task: str


class SkillRecord(BaseModel):
    employee: str
    skill: str


class SearchResult(BaseModel):
    count: int
    results: list


class ScanSummary(BaseModel):
    employees: int
    task_assignments: int
    skill_assignments: int
    file_count: int


def clean_name(value: str) -> str:
    if not value:
        return ""
    value = value.replace("\xa0", " ").replace("–", "-").replace("—", " - ")
    value = re.sub(r"[\u2022\u25aa\u25cf\uf0b7]+", " ", value)
    value = re.sub(r"\s+", " ", value)
    return value.strip(" |:-")


def resolve_report_path(raw_path: str) -> Path | None:
    if not raw_path:
        return None

    parsed = urlparse(raw_path)
    if parsed.scheme in {"http", "https"}:
        if "/uploads/" not in parsed.path:
            return None
        candidate = unquote(parsed.path)
    else:
        candidate = unquote(parsed.path if parsed.scheme == "file" else raw_path)

    path = Path(candidate)
    if path.is_absolute():
        return path

    normalized = candidate.replace("\\", "/").lstrip("/")
    if normalized.startswith("src/uploads/"):
        normalized = normalized[len("src/uploads/") :]
    if normalized.startswith("uploads/"):
        normalized = normalized[len("uploads/") :]
    return FILE_BASE_DIR / normalized


def field_values(document: dict, dotted_path: str) -> list[str]:
    values = [document]
    for part in dotted_path.split("."):
        next_values = []
        for value in values:
            if isinstance(value, list):
                next_values.extend(item.get(part) for item in value if isinstance(item, dict) and part in item)
            elif isinstance(value, dict) and part in value:
                next_values.append(value[part])
        values = next_values

    flattened = []
    for value in values:
        if isinstance(value, list):
            flattened.extend(item for item in value if isinstance(item, str))
        elif isinstance(value, str):
            flattened.append(value)
    return flattened


def read_docx_lines(path: Path) -> list[str]:
    with ZipFile(path) as zf:
        xml = zf.read("word/document.xml")
    root = ET.fromstring(xml)
    lines = []
    for para in root.findall(".//w:p", NS):
        parts = [node.text for node in para.findall(".//w:t", NS) if node.text]
        line = "".join(parts).strip()
        if line:
            lines.append(clean_name(line))
    return lines


def read_pdf_lines(path: Path) -> list[str]:
    if PdfReader:
        try:
            reader = PdfReader(str(path))
            lines = []
            for page in reader.pages:
                page_text = page.extract_text() or ""
                lines.extend(clean_name(line) for line in page_text.splitlines() if clean_name(line))
            if lines:
                return lines
        except Exception as exc:
            log.debug("PyPDF2 failed for %s: %s", path, exc)

    if fitz:
        doc = fitz.open(str(path))
        try:
            lines = []
            for page in doc:
                lines.extend(clean_name(line) for line in page.get_text().splitlines() if clean_name(line))
            return lines
        finally:
            doc.close()
    return []


def read_excel_lines(path: Path) -> list[str]:
    if not load_workbook:
        return []
    lines = []
    try:
        wb = load_workbook(path, data_only=True)
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                for cell in row:
                    if cell is not None:
                        cell_str = str(cell).strip()
                        if cell_str:
                            lines.append(clean_name(cell_str))
    except Exception as exc:
        log.debug("Excel extraction failed for %s: %s", path, exc)
    return lines


def read_pptx_lines(path: Path) -> list[str]:
    if not Presentation:
        return []
    lines = []
    try:
        prs = Presentation(str(path))
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    lines.append(clean_name(shape.text))
    except Exception as exc:
        log.debug("PPTX extraction failed for %s: %s", path, exc)
    return lines


def read_file_lines(path: Path) -> list[str]:
    ext = path.suffix.lower()
    if ext == ".docx":
        return read_docx_lines(path)
    if ext == ".pdf":
        return read_pdf_lines(path)
    if ext in (".xlsx", ".xlsm", ".xltx", ".xltm"):
        return read_excel_lines(path)
    if ext == ".pptx":
        return read_pptx_lines(path)
    return []


def extract_employee_name(lines: list[str], path: Path) -> str:
    for idx, line in enumerate(lines):
        if line.lower() == "employee name" and idx + 1 < len(lines):
            name = clean_name(lines[idx + 1])
            if len(name) > 1:
                return name

    parts = re.split(r"[_\-\s]+", path.stem)
    for idx, part in enumerate(parts):
        if part.lower() in {"report", "be", "fe", "ml", "data"} and idx + 2 < len(parts):
            return clean_name(f"{parts[idx + 1]} {parts[idx + 2]}")

    return clean_name(lines[0][:50]) if lines else path.stem


def extract_tasks(lines: list[str], path: Path) -> list[str]:
    tasks = []
    for idx, line in enumerate(lines):
        if line.lower() == "description" and idx > 0:
            title = clean_name(re.sub(r"^[^A-Za-z0-9]+", "", lines[idx - 1]))
            if 3 <= len(title) <= 120:
                tasks.append(title)

    if tasks:
        return sorted(set(tasks))

    title = clean_name(lines[0] if lines else path.stem)
    return [title] if title else []


def extract_skills(lines: list[str]) -> list[str]:
    skills: set[str] = set()
    for idx, line in enumerate(lines):
        lower = line.lower()
        if lower.startswith("technologies / tools") or lower.startswith("tools used"):
            if idx + 1 < len(lines):
                for piece in re.split(r",|/| and ", lines[idx + 1]):
                    piece = clean_name(piece)
                    if 2 <= len(piece) <= 50:
                        skills.add(piece)
    return sorted(skills)


def scan_summary() -> ScanSummary:
    employees = {r["employee"] for r in _store["tasks"]} | {r["employee"] for r in _store["skills"]}
    return ScanSummary(
        employees=len(employees),
        task_assignments=len(_store["tasks"]),
        skill_assignments=len(_store["skills"]),
        file_count=_store["file_count"],
    )


async def _run_scan() -> None:
    if mongo_client is None:
        raise RuntimeError("MongoDB client is not initialized")

    collection = mongo_client[MONGO_DB][MONGO_COLLECTION]
    cursor = collection.find({}, {FILE_PATH_FIELD: 1, "_id": 0})
    docs = await cursor.to_list(length=None)

    if not docs:
        log.warning("No documents found in %s.%s", MONGO_DB, MONGO_COLLECTION)
        _store.update({"tasks": [], "skills": [], "scanned": True, "file_count": 0})
        return

    task_records, skill_records = [], []
    processed = 0

    for doc in docs:
        raw_paths = field_values(doc, FILE_PATH_FIELD)
        if not raw_paths:
            log.warning("Document missing '%s' field: %s", FILE_PATH_FIELD, doc)
            continue

        for raw_path in raw_paths:
            path = resolve_report_path(raw_path)
            if path is None:
                log.warning("Unsupported remote or empty path in '%s': %s", FILE_PATH_FIELD, raw_path)
                continue
            if not path.exists():
                log.warning("File not found on disk: %s", path)
                continue

            lines = read_file_lines(path)
            if not lines:
                log.warning("No text extracted from %s", path.name)
                continue

            employee = extract_employee_name(lines, path)
            if not employee:
                continue

            task_records.extend({"employee": employee, "task": task} for task in extract_tasks(lines, path))
            skill_records.extend({"employee": employee, "skill": skill} for skill in extract_skills(lines))
            processed += 1
            log.info("Processed %s -> %s", path.name, employee)

    tasks_df = pd.DataFrame(task_records).drop_duplicates() if task_records else pd.DataFrame(columns=["employee", "task"])
    skills_df = pd.DataFrame(skill_records).drop_duplicates() if skill_records else pd.DataFrame(columns=["employee", "skill"])

    _store["tasks"] = tasks_df.to_dict("records")
    _store["skills"] = skills_df.to_dict("records")
    _store["scanned"] = True
    _store["file_count"] = processed
    log.info("Scan complete: %s files, %s tasks, %s skills", processed, len(_store["tasks"]), len(_store["skills"]))


@asynccontextmanager
async def lifespan(app: FastAPI):
    global mongo_client
    mongo_client = AsyncIOMotorClient(MONGO_URI)
    log.info("Connected to MongoDB: %s / db=%s / collection=%s", MONGO_URI, MONGO_DB, MONGO_COLLECTION)
    try:
        await _run_scan()
    except Exception as exc:
        log.warning("Initial report scan failed: %s", exc)
    yield
    mongo_client.close()
    log.info("MongoDB connection closed")


app = FastAPI(
    title="Employee Skill & Task Finder API",
    description="Reads report file paths from MongoDB and extracts employee tasks and skills.",
    version="3.0.0",
    lifespan=lifespan,
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)


@app.get("/healthz")
def healthz():
    return {"status": "ok", "scanned": _store["scanned"], "file_count": _store["file_count"]}


@app.get("/")
def root():
    return {
        "status": "ok",
        "scanned": _store["scanned"],
        "file_count": _store["file_count"],
        "mongo": f"{MONGO_URI} / {MONGO_DB}.{MONGO_COLLECTION}",
        "file_path_field": FILE_PATH_FIELD,
        "file_base_dir": str(FILE_BASE_DIR),
    }


@app.post("/rescan", response_model=ScanSummary)
async def rescan():
    await _run_scan()
    return scan_summary()


@app.get("/status", response_model=ScanSummary)
def status():
    return scan_summary()


@app.get("/skills", response_model=SearchResult)
def search_skills(q: Optional[str] = Query(None, description="Partial skill name")):
    records = _store["skills"]
    if q:
        records = [r for r in records if q.lower() in r["skill"].lower()]
    return SearchResult(count=len(records), results=records)


@app.get("/tasks", response_model=SearchResult)
def search_tasks(q: Optional[str] = Query(None, description="Partial task name")):
    records = _store["tasks"]
    if q:
        records = [r for r in records if q.lower() in r["task"].lower()]
    return SearchResult(count=len(records), results=records)


@app.get("/employees")
def list_employees():
    names = sorted({r["employee"] for r in _store["tasks"]} | {r["employee"] for r in _store["skills"]})
    return {"count": len(names), "employees": names}


@app.get("/search/skills")
def search_employees_by_skill(q: str = Query(..., min_length=1, description="Skill keyword")):
    matched = [r for r in _store["skills"] if q.lower() in r["skill"].lower()]
    if not matched:
        raise HTTPException(status_code=404, detail=f"No employees found with skill matching '{q}'.")

    grouped: dict[str, list[str]] = {}
    for record in matched:
        grouped.setdefault(record["employee"], []).append(record["skill"])

    return {
        "query": q,
        "count": len(matched),
        "employees": [
            {"employee": employee, "matched_skills": sorted(set(skills))}
            for employee, skills in sorted(grouped.items())
        ],
    }


@app.get("/search/tasks")
def search_employees_by_task(q: str = Query(..., min_length=1, description="Task keyword")):
    matched = [r for r in _store["tasks"] if q.lower() in r["task"].lower()]
    if not matched:
        raise HTTPException(status_code=404, detail=f"No employees found with task matching '{q}'.")

    grouped: dict[str, list[str]] = {}
    for record in matched:
        grouped.setdefault(record["employee"], []).append(record["task"])

    return {
        "query": q,
        "count": len(matched),
        "employees": [
            {"employee": employee, "matched_tasks": sorted(set(tasks))}
            for employee, tasks in sorted(grouped.items())
        ],
    }
